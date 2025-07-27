import os
from fastapi import UploadFile
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

load_dotenv()
UPLOAD_DIR = os.getenv("UPLOAD_DIR", "uploads")
REWORKED_DIR = os.path.join(UPLOAD_DIR, "reworked")
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(REWORKED_DIR, exist_ok=True)

async def save_file(file: UploadFile) -> str:
    original_path = os.path.join(UPLOAD_DIR, file.filename)
    # сохраняем исходный файл
    with open(original_path, "wb") as f:
        content = await file.read()
        f.write(content)
    # редизайн и сохранение
    new_filename = f"redesigned_{file.filename}"
    redesigned_path = os.path.join(REWORKED_DIR, new_filename)
    redesign_presentation(original_path, redesigned_path)
    return f"/{REWORKED_DIR}/{new_filename}"

def redesign_presentation(input_path: str, output_path: str):
    prs = Presentation(input_path)
    # Цветовая палитра "древний пергамент"
    bg_color = RGBColor(245, 240, 230)  # цвет старой бумаги
    text_color = RGBColor(64, 42, 32)   # темно-коричневый текст
    accent_color = RGBColor(101, 67, 33)  # коричневый для заголовков
    secondary_color = RGBColor(139, 90, 43)  # светло-коричневый для акцентов
    
    for i, slide in enumerate(prs.slides):
        # Устанавливаем фон
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        
        # Удаляем все существующие линии и соединители
        remove_existing_lines(slide)
        
        title_found = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Определяем тип текста - сначала проверяем placeholder, потом позицию
                is_title = False
                
                # Проверяем является ли это title placeholder
                try:
                    if (hasattr(shape, "placeholder_format") and 
                        shape.placeholder_format is not None and
                        shape.placeholder_format.idx == 0 and
                        not title_found):
                        is_title = True
                except:
                    # Если ошибка с placeholder, определяем по позиции и размеру
                    if (shape.top < Inches(2) and 
                        shape.width > Inches(6) and 
                        not title_found):
                        is_title = True
                
                if is_title:
                    title_found = True
                    format_title(shape, accent_color)
                else:
                    format_body_text(shape, text_color)
                
                # Настраиваем автоподгонку текста
                setup_text_autofit(shape)
                
                # Добавляем отступы для лучшей читаемости
                try:
                    shape.text_frame.margin_left = Inches(0.3)
                    shape.text_frame.margin_right = Inches(0.3)
                    shape.text_frame.margin_top = Inches(0.15)
                    shape.text_frame.margin_bottom = Inches(0.15)
                except:
                    pass  # игнорируем ошибки настройки отступов
        
        title_shape = ""

        # Добавляем декоративную линию под заголовком (только если есть заголовок)
        if title_found and i > 0:  # не на титульном слайде
            add_decorative_line(slide, secondary_color , title_shape)
    
    prs.save(output_path)

def remove_existing_lines(slide):
    """Удаляет все существующие линии и соединители со слайда"""
    shapes_to_remove = []
    for shape in slide.shapes:
        try:
            # Проверяем различные типы линий
            if (shape.shape_type == MSO_SHAPE_TYPE.LINE or 
                shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or
                (hasattr(shape, 'line') and hasattr(shape.line, 'width') and shape.line.width > Pt(0.5))):
                # Дополнительная проверка - если у объекта есть line но нет text_frame
                if not shape.has_text_frame:
                    shapes_to_remove.append(shape)
        except:
            continue
    
    # Удаляем найденные линии
    for shape in shapes_to_remove:
        try:
            slide.shapes._spTree.remove(shape._element)
        except:
            pass

def setup_text_autofit(shape):
    """Настраивает автоподгонку текста чтобы он не выходил за границы"""
    try:
        text_frame = shape.text_frame
        # Включаем автоподгонку текста
        text_frame.auto_size = None  # отключаем автоизменение размера
        text_frame.word_wrap = True  # включаем перенос слов
        
        # Устанавливаем вертикальное выравнивание по верху
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
    except Exception as e:
        pass  # игнорируем ошибки настройки

def format_title(shape, accent_color):
    """Форматирование заголовков с проверкой переполнения"""
    try:
        # Рассчитываем размер шрифта в зависимости от длины текста
        total_text_length = sum(len(para.text) for para in shape.text_frame.paragraphs)
        
        if total_text_length > 60:
            title_size = Pt(24)  # меньший размер для длинных заголовков
        elif total_text_length > 30:
            title_size = Pt(28)
        else:
            title_size = Pt(32)
            
        for para in shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            # Устанавливаем межстрочный интервал
            para.space_before = Pt(0)
            para.space_after = Pt(8)
            
            for run in para.runs:
                run.font.name = 'Times New Roman'  # классический шрифт для пергамента
                run.font.size = title_size
                run.font.bold = True
                run.font.italic = False
                run.font.color.rgb = accent_color
                
    except Exception as e:
        pass

def format_body_text(shape, text_color):
    """Форматирование основного текста с проверкой переполнения"""
    try:
        # Рассчитываем размер шрифта в зависимости от количества текста
        total_text_length = sum(len(para.text) for para in shape.text_frame.paragraphs)
        
        if total_text_length > 400:
            body_size = Pt(14)  # меньший размер для очень длинного текста
        elif total_text_length > 200:
            body_size = Pt(16)
        else:
            body_size = Pt(18)
            
        for para in shape.text_frame.paragraphs:
            # Автоматическое выравнивание
            if len(para.text.strip()) < 80:
                para.alignment = PP_ALIGN.LEFT
            else:
                para.alignment = PP_ALIGN.JUSTIFY
                
            # Настройка интервалов
            para.space_before = Pt(3)
            para.space_after = Pt(6)
            para.line_spacing = 1.2  # межстрочный интервал
                
            for run in para.runs:
                run.font.name = 'Georgia'  # читаемый шрифт для основного текста
                run.font.size = body_size
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = text_color
                
    except Exception as e:
        pass

def add_decorative_line(slide, secondary_color, title_shape):
    """Добавляет декоративную линию под заголовком в стиле пергамента"""
    try:
        # Рассчитываем позицию линии на основе реального заголовка
        title_bottom = title_shape.top + title_shape.height
        line_y = title_bottom + Inches(0.1)  # небольшой отступ под заголовком
        
        # Используем границы заголовка для определения длины линии
        line_left = title_shape.left
        line_right = title_shape.left + title_shape.width
        
        # Создаем более толстую линию в стиле старинного документа
        line = slide.shapes.add_connector(
            1,  # прямая линия
            line_left, line_y,      # начальная точка
            line_right, line_y      # конечная точка
        )
        line.line.color.rgb = secondary_color
        line.line.width = Pt(3)  # более толстая линия
        
        # Добавляем вторую более тонкую линию для эффекта "двойной линии"
        thin_line = slide.shapes.add_connector(
            1,
            line_left, line_y + Inches(0.05),  # чуть ниже первой линии
            line_right, line_y + Inches(0.05)
        )
        thin_line.line.color.rgb = secondary_color
        thin_line.line.width = Pt(1)
        
    except Exception as e:
        pass  # игнорируем ошибки добавления линии